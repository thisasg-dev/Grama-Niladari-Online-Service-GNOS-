"""
GNOS RAG Bot - v6.0 (GEMINI PRIMARY + GROQ FALLBACK)
═══════════════════════════════════════════════════════════════════
MAJOR CHANGES from v5.7:

1. GEMINI IS NOW PRIMARY (was fallback before)
   WHY: gemini-2.5-flash has far better Sinhala & Tamil support
        than LLaMA. Google trained on much more South Asian data.
        Also 1M tokens/day free vs 100k for Groq.

2. GROQ IS NOW FALLBACK (was primary before)
   WHY: Still very fast for English. Good safety net.
        Used when Gemini hits rate limits.

3. FIXED GEMINI TOOL CALL BUG
   BUG: Bot was saying "I will search govlk_search(...)" as TEXT
        instead of actually calling the tool.
   FIX: Use response.function_calls (official SDK method) instead
        of manually checking parts[].function_call
        Also disabled automatic_function_calling so we control it.

4. GEMINI ANSWERS SINHALA/TAMIL NATIVELY
   OLD: LLM answers in English → Google Translate → Sinhala/Tamil
        (translation errors creep in)
   NEW: Gemini answers directly in Sinhala/Tamil
        (much more accurate, natural phrasing)
   NOTE: Groq still uses Google Translate (LLaMA is English-first)

PROVIDER CHAIN:
  Gemini Key_1  (1M tokens/day)   ← PRIMARY
  Gemini Key_2  (1M tokens/day)   ← Gemini fallback (optional)
  Groq Key_1    (100k tokens/day) ← Secondary fallback
  Groq Key_2    (100k tokens/day) ← Tertiary fallback
  Groq Key_3    (100k tokens/day) ← Last resort
  ─────────────────────────────────────────────────
  TOTAL: ~2.3M tokens/day FREE

.env file:
  GEMINI_API_KEY_1=AIzaSy...   ← primary Gemini (aistudio.google.com)
  GEMINI_API_KEY_2=AIzaSy...   ← optional second Gemini account
  GROQ_API_KEY_1=gsk_...       ← fallback Groq account 1
  GROQ_API_KEY_2=gsk_...       ← fallback Groq account 2
  GROQ_API_KEY_3=gsk_...       ← fallback Groq account 3
  HF_TOKEN=hf_...              ← HuggingFace (optional)

Install:
  pip install google-genai groq ddgs sentence-transformers
  pip install faiss-cpu langchain-community langdetect
  pip install deep-translator python-docx python-pptx pypdf
═══════════════════════════════════════════════════════════════════
"""

import os
import re
import json
import time
import pickle
import hashlib
import datetime

import faiss
import numpy as np
from dotenv import load_dotenv
from langdetect import detect
from deep_translator import GoogleTranslator
from sentence_transformers import SentenceTransformer
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from docx import Document
from pptx import Presentation
from huggingface_hub import login

# ── DDGS import ───────────────────────────────────────────────────
try:
    from ddgs import DDGS
    print("  Using ddgs")
except ImportError:
    try:
        from duckduckgo_search import DDGS
        print("  Using duckduckgo_search (run: pip install ddgs)")
    except ImportError:
        raise ImportError("Run: pip install ddgs")

# ── Gemini import ─────────────────────────────────────────────────
try:
    from google import genai as google_genai
    from google.genai import types as gtypes
    GEMINI_AVAILABLE = True
    print("  google-genai SDK: loaded")
except ImportError:
    GEMINI_AVAILABLE = False
    print("  ⚠ Run: pip install google-genai")

# ── Groq import ───────────────────────────────────────────────────
try:
    from groq import Groq
    GROQ_AVAILABLE = True
    print("  groq SDK: loaded")
except ImportError:
    GROQ_AVAILABLE = False
    print("  ⚠ Groq not installed (optional fallback)")

load_dotenv()

# ── Safe HuggingFace login ────────────────────────────────────────
try:
    hf_token = os.getenv("HF_TOKEN")
    if hf_token:
        login(token=hf_token)
        print("  HuggingFace: logged in")
    else:
        print("  HuggingFace: skipped (no HF_TOKEN)")
except Exception as e:
    print(f"  HuggingFace: skipped ({e})")


# ═══════════════════════════════════════════════════════════════════
# GEMINI CLIENT POOL  (PRIMARY provider)
# ═══════════════════════════════════════════════════════════════════
GEMINI_PRIMARY_MODEL  = "gemini-2.5-flash"       # best: Sinhala/Tamil/English
GEMINI_FALLBACK_MODEL = "gemini-2.0-flash"        # lighter fallback

GEMINI_CLIENTS = []   # list of (client, label)
_gemini_idx    = 0

if GEMINI_AVAILABLE:
    # Load GEMINI_API_KEY_1, _2, etc. OR legacy GEMINI_API_KEY
    gemini_keys = []
    for i in range(1, 4):
        k = os.getenv(f"GEMINI_API_KEY_{i}")
        if k and k.strip():
            gemini_keys.append((k.strip(), f"Gemini_Key_{i}"))
    if not gemini_keys:
        k = os.getenv("GEMINI_API_KEY")
        if k and k.strip():
            gemini_keys.append((k.strip(), "Gemini_Key_1"))

    for key, label in gemini_keys:
        try:
            c = google_genai.Client(api_key=key)
            GEMINI_CLIENTS.append((c, label))
            print(f"  {label}: loaded (~1M tokens/day)")
        except Exception as e:
            print(f"  {label}: failed — {e}")

def get_gemini_client():
    return GEMINI_CLIENTS[_gemini_idx][0]

def rotate_gemini() -> bool:
    global _gemini_idx
    if _gemini_idx + 1 < len(GEMINI_CLIENTS):
        _gemini_idx += 1
        print(f"  🔄 Rotated to {GEMINI_CLIENTS[_gemini_idx][1]}")
        return True
    return False

def reset_gemini():
    global _gemini_idx
    _gemini_idx = 0


# ═══════════════════════════════════════════════════════════════════
# GROQ CLIENT POOL  (FALLBACK provider)
# ═══════════════════════════════════════════════════════════════════
GROQ_MODEL      = "llama-3.3-70b-versatile"
GROQ_MODEL_FAST = "llama-3.1-8b-instant"

GROQ_CLIENTS = []   # list of (client, label)
_groq_idx    = 0

if GROQ_AVAILABLE:
    groq_keys = []
    for i in range(1, 6):
        k = os.getenv(f"GROQ_API_KEY_{i}")
        if k and k.strip():
            groq_keys.append((k.strip(), f"Groq_Key_{i}"))
    if not groq_keys:
        k = os.getenv("GROQ_API_KEY")
        if k and k.strip():
            groq_keys.append((k.strip(), "Groq_Key_1"))

    for key, label in groq_keys:
        try:
            c = Groq(api_key=key)
            GROQ_CLIENTS.append((c, label))
            print(f"  {label}: loaded (~100k tokens/day)")
        except Exception as e:
            print(f"  {label}: failed — {e}")

def get_groq_client():
    return GROQ_CLIENTS[_groq_idx][0]

def rotate_groq() -> bool:
    global _groq_idx
    if _groq_idx + 1 < len(GROQ_CLIENTS):
        _groq_idx += 1
        print(f"  🔄 Rotated to {GROQ_CLIENTS[_groq_idx][1]}")
        return True
    return False

def reset_groq():
    global _groq_idx
    _groq_idx = 0

# Print total capacity
gemini_cap = len(GEMINI_CLIENTS) * 1000
groq_cap   = len(GROQ_CLIENTS) * 100
print(f"  Total daily capacity: ~{gemini_cap + groq_cap}k tokens/day")

if not GEMINI_CLIENTS and not GROQ_CLIENTS:
    raise ValueError("No API keys found. Add GEMINI_API_KEY_1 or GROQ_API_KEY_1 to .env")


# ═══════════════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════════════
MAX_TOOL_ROUNDS  = 5
MAX_SEARCH_CHARS = 6000
CACHE_TTL        = 300   # 5 minutes for live question cache
FAISS_FILE       = "faiss_index.index"
DOCS_FILE        = "docs.pkl"


# ═══════════════════════════════════════════════════════════════════
# DATE HELPERS
# ═══════════════════════════════════════════════════════════════════
def get_today() -> str:
    return datetime.datetime.now().strftime("%A, %B %d, %Y")

def get_year() -> int:
    return datetime.datetime.now().year


# ═══════════════════════════════════════════════════════════════════
# ANSWER CACHE  (live questions only, 5-min TTL)
# ═══════════════════════════════════════════════════════════════════
_answer_cache: dict = {}

def _ckey(q: str) -> str:
    return hashlib.md5(q.lower().strip().encode()).hexdigest()

def get_cached(q: str):
    item = _answer_cache.get(_ckey(q))
    if item:
        ans, ts = item
        age = time.time() - ts
        if age < CACHE_TTL:
            print(f"  ⚡ Cache hit ({int(age)}s old)")
            return ans
        del _answer_cache[_ckey(q)]
    return None

def set_cached(q: str, ans: str):
    _answer_cache[_ckey(q)] = (ans, time.time())
    print(f"  💾 Cached ({CACHE_TTL}s TTL)")


# ═══════════════════════════════════════════════════════════════════
# KEYWORD DETECTORS
# ═══════════════════════════════════════════════════════════════════
FORCE_WEB_KEYWORDS = {
    "upgrade","refresh","revise","modify","amend","renew","modernize",
    "enhance","edit","patch","improve","refine","adjust","rework",
    "overhaul","fix","synchronize","recondition","renovate","reform",
    "recent","current","up-to-date","newest","fresh","modern",
    "contemporary","just-released","present","trending","newly-issued",
    "novel","brand-new","innovative","cutting-edge","state-of-the-art",
    "newfangled","advanced","pioneering","mint","untouched","unfamiliar",
    "current-day","today","tomorrow","up-to-the-minute","hot","breaking",
    "latest","who","whom","person","people",
    "president","prime minister","minister","cabinet","cabinets",
    "parliament","government","mp","senator","governor","ambassador",
    "secretary","chairman","director","leader","appointed","elected",
    "opposition","law","laws","gazette","gazettes","regulation",
    "regulations","policy","policies","act","bill","ordinance",
    "amendment","legislation","circular","legal","court","judgment",
    "constitution","statutory","world","global","international",
    "country","countries","land","nation","sri lanka","srilanka",
    "lanka","colombo","most","biggest","largest","smallest","best",
    "worst","richest","poorest","fastest","slowest","highest","lowest",
    "top","list","value","price","prices","rate","cost","budget","tax",
    "economy","inflation","exchange","dollar","rupee","fuel","future",
    "requirements","upcoming","2020","2021","2022","2023","2024",
    "2025","2026","2027",
}

GOVLK_KEYWORDS = {
    "president","prime minister","minister","cabinet","opposition",
    "parliament","mp","parliamentarian","speaker","appointed","elected",
    "sworn","law","laws","act","bill","gazette","gazettes","regulation",
    "regulations","ordinance","amendment","legislation","circular",
    "legal notice","statutory","constitution","policy","policies",
    "government policy","cabinet decision","government news",
    "official announcement","press release","government statement",
}

GOVLK_SOURCES = [
    ("site:presidentsoffice.gov.lk",  "President's Office"),
    ("site:parliament.lk",            "Sri Lanka Parliament"),
    ("site:documents.gov.lk",         "Official Gazette Documents"),
    ("site:cabinetoffice.gov.lk",     "Cabinet Office"),
    ("site:treasury.gov.lk",          "Ministry of Finance"),
    ("site:dgi.gov.lk",               "Dept of Government Information"),
    ("site:gov.lk",                   "Sri Lanka Government Portal"),
]

def must_use_web(q: str) -> bool:
    ql = q.lower()
    for kw in FORCE_WEB_KEYWORDS:
        if kw in ql:
            print(f"  🔴 FORCE WEB: '{kw}'")
            return True
    return False

def is_govlk_topic(q: str) -> bool:
    ql    = q.lower()
    is_sl = any(k in ql for k in ["sri lanka","srilanka","lanka","colombo","lk"])
    is_gv = any(k in ql for k in GOVLK_KEYWORDS)
    if is_sl and is_gv: return True
    strong = {"president","prime minister","cabinet","parliament",
              "gazette","opposition","constitution"}
    return any(k in ql for k in strong)

def is_live(q: str) -> bool:
    return is_govlk_topic(q) or must_use_web(q)


# ═══════════════════════════════════════════════════════════════════
# LANGUAGE HELPERS
# ═══════════════════════════════════════════════════════════════════
def contains_sinhala(text: str) -> bool:
    return any('\u0D80' <= c <= '\u0DFF' for c in text)

def contains_tamil(text: str) -> bool:
    return any('\u0B80' <= c <= '\u0BFF' for c in text)

def detect_language(text: str) -> str:
    if contains_sinhala(text): return "Sinhala"
    if contains_tamil(text):   return "Tamil"
    try:
        code = detect(text)
        if code == "si": return "Sinhala"
        if code == "ta": return "Tamil"
    except: pass
    return "English"

def detect_preferred_language(q: str) -> str:
    ql = q.lower()
    if any(k in ql for k in ["sinhala","සිංහල","in si","sinhala language","sinhalen"]):
        return "Sinhala"
    if any(k in ql for k in ["tamil","தமிழ்","in ta","tamil language"]):
        return "Tamil"
    if any(k in ql for k in ["english","in english","english language"]):
        return "English"
    return detect_language(q)

def translate_to_english(text: str) -> str:
    try: return GoogleTranslator(source='auto', target='en').translate(text)
    except: return text

def translate_answer(text: str, lang: str) -> str:
    code = 'si' if lang == "Sinhala" else 'ta'
    try:
        lines = text.split('\n')
        chunks, cur, cur_len = [], [], 0
        for line in lines:
            ll = len(line)
            if cur_len + ll + 1 > 800 and cur:
                chunks.append('\n'.join(cur)); cur, cur_len = [], 0
            cur.append(line); cur_len += ll + 1
        if cur: chunks.append('\n'.join(cur))
        out = []
        for chunk in chunks:
            chunk = chunk.strip()
            if not chunk: out.append(''); continue
            try:
                t = GoogleTranslator(source='en', target=code).translate(chunk)
                out.append(t or chunk)
            except: out.append(chunk)
        return '\n'.join(out)
    except: return text

def smart_translate(text: str, lang: str, used_gemini: bool = False) -> str:
    """
    Gemini can answer in Sinhala/Tamil directly — no translation needed.
    Groq (LLaMA) needs Google Translate for Sinhala/Tamil.
    """
    if lang == "English": return text
    if used_gemini:
        # Gemini answered natively — check if already in correct script
        if lang == "Sinhala" and contains_sinhala(text): return text
        if lang == "Tamil"   and contains_tamil(text):   return text
        # If not native, translate as fallback
    # Groq path or Gemini gave English despite native request
    if lang == "Sinhala": return translate_answer(text, "Sinhala")
    if lang == "Tamil":
        if contains_tamil(text): return text
        return translate_answer(text, "Tamil")
    return text

def extract_wait(err: str) -> int:
    try:
        m = re.search(r'try again in\s+(?:(\d+)m)?(?:([\d.]+)s)?', err, re.I)
        if m: return min(int(int(m.group(1) or 0)*60 + float(m.group(2) or 0)) + 2, 120)
    except: pass
    return 30


# ═══════════════════════════════════════════════════════════════════
# TOOL IMPLEMENTATIONS
# ═══════════════════════════════════════════════════════════════════
def tool_web_search(query: str) -> str:
    print(f"  🌐 web_search: {repr(query)}")
    results, seen = [], set()
    try:
        time.sleep(0.5)
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=8):
                body  = r.get("body",  "").strip()
                title = r.get("title", "").strip()
                href  = r.get("href",  "")
                if body and body not in seen:
                    seen.add(body)
                    e = f"[{title}]\n{body}"
                    if href: e += f"\n(URL: {href})"
                    results.append(e)
    except Exception as ex:
        return f"Web search failed: {ex}"
    if not results: return f"No results for: {query}"
    c = "\n\n---\n\n".join(results)
    if len(c) > MAX_SEARCH_CHARS: c = c[:MAX_SEARCH_CHARS] + "\n[truncated]"
    print(f"  ✓ {len(results)} results")
    return c

def tool_govlk_search(query: str) -> str:
    print(f"  🏛️  govlk_search: {repr(query)}")
    year, all_r, seen = get_year(), [], set()
    for sf, sn in GOVLK_SOURCES:
        try:
            time.sleep(0.4)
            with DDGS() as ddgs:
                for r in ddgs.text(f"{query} {year} {sf}", max_results=4):
                    body  = r.get("body",  "").strip()
                    title = r.get("title", "").strip()
                    href  = r.get("href",  "")
                    if body and body not in seen:
                        seen.add(body)
                        e = f"[{sn} — {title}]\n{body}"
                        if href: e += f"\n(URL: {href})"
                        all_r.append(e)
        except Exception as ex:
            print(f"  ⚠ {sn}: {ex}")
    if all_r:
        c = "\n\n---\n\n".join(all_r)
        if len(c) > MAX_SEARCH_CHARS: c = c[:MAX_SEARCH_CHARS] + "\n[truncated]"
        return f"[Source: Official Sri Lanka Government — gov.lk]\n\n{c}"
    print(f"  ⚠ gov.lk empty → DuckDuckGo")
    return tool_web_search(query)

def tool_search_documents(query: str) -> str:
    print(f"  📄 search_documents: {repr(query)}")
    try:
        qe   = translate_to_english(query)
        qemb = np.array(embedding_model.encode([qe])).astype("float32")
        dist, idxs = index.search(qemb, 5)
        score = dist[0][0]
        print(f"  FAISS score: {score:.4f}")
        if score > 0.8: return "No relevant documents found."
        results = []
        for i in idxs[0][:3]:
            t = docs[i]["text"].strip()
            s = os.path.basename(docs[i]["source"])
            if t: results.append(f"[From: {s}]\n{t}")
        return "\n\n---\n\n".join(results) if results else "No relevant documents found."
    except Exception as ex:
        return f"Document search error: {ex}"

def dispatch(name: str, args: dict) -> str:
    if name == "govlk_search":      return tool_govlk_search(args.get("query",""))
    elif name == "web_search":      return tool_web_search(args.get("query",""))
    elif name == "search_documents": return tool_search_documents(args.get("query",""))
    return f"Unknown tool: {name}"


# ═══════════════════════════════════════════════════════════════════
# GEMINI TOOL DEFINITIONS
# ═══════════════════════════════════════════════════════════════════
def make_gemini_tools():
    if not GEMINI_AVAILABLE: return []
    return [gtypes.Tool(function_declarations=[
        gtypes.FunctionDeclaration(
            name="govlk_search",
            description=(
                "Search official Sri Lanka government websites (gov.lk). "
                "USE FIRST for: President, Prime Minister, Ministers, Cabinet, "
                "Opposition, Parliament, laws, Gazettes, regulations, policies, "
                "government news. Sources include presidentsoffice.gov.lk, "
                "parliament.lk, documents.gov.lk, cabinetoffice.gov.lk"
            ),
            parameters=gtypes.Schema(
                type=gtypes.Type.OBJECT,
                properties={"query": gtypes.Schema(
                    type=gtypes.Type.STRING,
                    description="Search query. Always include year. E.g. 'President Sri Lanka 2026'"
                )},
                required=["query"]
            )
        ),
        gtypes.FunctionDeclaration(
            name="web_search",
            description=(
                "Search internet for current real-time information: "
                "news, prices, world events, current status of anything. "
                "For Sri Lanka government topics use govlk_search first."
            ),
            parameters=gtypes.Schema(
                type=gtypes.Type.OBJECT,
                properties={"query": gtypes.Schema(
                    type=gtypes.Type.STRING,
                    description="Search query with year for current info."
                )},
                required=["query"]
            )
        ),
        gtypes.FunctionDeclaration(
            name="search_documents",
            description=(
                "Search local PDF/document database for: "
                "application procedures, required documents, "
                "Grama Niladhari services, certificate processes."
            ),
            parameters=gtypes.Schema(
                type=gtypes.Type.OBJECT,
                properties={"query": gtypes.Schema(
                    type=gtypes.Type.STRING,
                    description="What to search in local documents."
                )},
                required=["query"]
            )
        ),
    ])]

GEMINI_TOOLS = make_gemini_tools()


# ═══════════════════════════════════════════════════════════════════
# GROQ TOOL DEFINITIONS
# ═══════════════════════════════════════════════════════════════════
GROQ_TOOLS = [
    {"type":"function","function":{"name":"govlk_search",
     "description":"Search official Sri Lanka government websites. USE FIRST for President, PM, Cabinet, Laws, Gazettes, Parliament, Policies.",
     "parameters":{"type":"object","properties":{"query":{"type":"string","description":"Include year. E.g. 'Prime Minister Sri Lanka 2026'"}},"required":["query"]}}},
    {"type":"function","function":{"name":"web_search",
     "description":"Search internet for current info, news, world events, prices.",
     "parameters":{"type":"object","properties":{"query":{"type":"string","description":"Search query with year."}},"required":["query"]}}},
    {"type":"function","function":{"name":"search_documents",
     "description":"Search local government PDF database for procedures and forms.",
     "parameters":{"type":"object","properties":{"query":{"type":"string","description":"Document search query."}},"required":["query"]}}},
]


# ═══════════════════════════════════════════════════════════════════
# SYSTEM PROMPT
# ═══════════════════════════════════════════════════════════════════
def build_prompt(answer_language: str,
                 force_web: bool   = False,
                 force_govlk: bool = False) -> str:
    today = get_today()
    year  = get_year()

    force_block = ""
    if force_govlk:
        force_block = (
            f"\n⚠ MANDATORY: Sri Lanka government question. "
            f"Call govlk_search FIRST with year {year}. "
            f"NEVER answer from training memory.\n"
        )
    elif force_web:
        force_block = (
            f"\n⚠ MANDATORY: Current info needed. "
            f"Call web_search FIRST with year {year}. "
            f"NEVER answer from training memory.\n"
        )

    return f"""You are GNOS Assistant — accurate, real-time government assistant for Sri Lanka.
Today: {today}.
{force_block}
TOOLS:
1. govlk_search — official Sri Lanka gov.lk sites (President, PM, Cabinet, Laws, Gazette, Parliament)
2. web_search   — internet search (current events, world news, prices, status)
3. search_documents — local PDF database (procedures, forms, GN services)

DECISION RULES:
- Sri Lanka government question → govlk_search FIRST, then web_search if incomplete
- Current events, news, prices, world → web_search FIRST
- Procedures, application forms → search_documents
- Simple math or well-known fact → answer directly (no search)

CRITICAL RULES:
- Today is {today}. ALWAYS search fresh for current facts.
- NEVER say "I will search" without actually calling the tool.
- NEVER answer current facts from training memory.
- ALWAYS give complete answers. Never say "not in context".
- If search is partial → use it AND fill rest from knowledge.

ANSWER FORMAT:
1. Direct answer first (1-2 sentences)
2. Details in numbered points
3. "Source: [where] — {today}"

LANGUAGE: Answer in {answer_language}.
IGNORE form field labels (Name:, Gender:, Date:, Line 1, etc.)
"""


# ═══════════════════════════════════════════════════════════════════
# GEMINI AGENTIC CALL  (FIXED tool detection)
# ═══════════════════════════════════════════════════════════════════
def gemini_agentic(question: str,
                   system_prompt: str,
                   history: list,
                   answer_language: str) -> str:
    """
    Full agentic loop using Gemini.
    
    FIX: Use response.function_calls (official SDK property)
         instead of manually iterating parts[].function_call
         This correctly detects ALL function calls including
         when Gemini batches multiple calls together.
    
    Native language: Gemini answers in Sinhala/Tamil directly.
    """
    # Build contents from history
    contents = []
    for msg in history:
        role    = msg.get("role","")
        content = msg.get("content","")
        if role in ("user", "assistant") and content:
            gemini_role = "model" if role == "assistant" else "user"
            contents.append(gtypes.Content(
                role=gemini_role,
                parts=[gtypes.Part(text=content)]
            ))

    # Add current question
    contents.append(gtypes.Content(
        role="user",
        parts=[gtypes.Part(text=question)]
    ))

    config = gtypes.GenerateContentConfig(
        system_instruction=system_prompt,
        tools=GEMINI_TOOLS,
        temperature=0.0,
        max_output_tokens=1500,
        # IMPORTANT: disable automatic calling so WE control the loop
        automatic_function_calling=gtypes.AutomaticFunctionCallingConfig(disable=True),
    )

    model = GEMINI_PRIMARY_MODEL

    for round_num in range(MAX_TOOL_ROUNDS):
        print(f"  🔮 Gemini round {round_num+1}/{MAX_TOOL_ROUNDS} | {model}")

        response = get_gemini_client().models.generate_content(
            model=model,
            contents=contents,
            config=config,
        )

        # ── FIXED: use response.function_calls ───────────────────
        # This is the correct official way to check for tool calls
        # Works even when Gemini batches multiple function calls
        if response.function_calls:
            # Add model's response to contents
            contents.append(response.candidates[0].content)

            # Execute each function call
            tool_response_parts = []
            for fc in response.function_calls:
                fn_name = fc.name
                fn_args = dict(fc.args) if fc.args else {}
                print(f"  🔮 Tool: {fn_name}({fn_args})")
                result = dispatch(fn_name, fn_args)
                tool_response_parts.append(
                    gtypes.Part.from_function_response(
                        name=fn_name,
                        response={"result": result}
                    )
                )

            # Add all tool responses as one user turn
            contents.append(gtypes.Content(
                role="user",
                parts=tool_response_parts
            ))

        else:
            # Final answer
            final = (response.text or "").strip()
            if not final:
                final = "I was unable to find a clear answer. Please try rephrasing."
            print(f"  🔮 Gemini answer ({len(final)} chars)")
            return final

    return "I searched but could not compose a complete answer. Please try a more specific question."


# ═══════════════════════════════════════════════════════════════════
# GROQ AGENTIC CALL  (fallback)
# ═══════════════════════════════════════════════════════════════════
def groq_agentic(question: str,
                 system_prompt: str,
                 history: list) -> str:
    """Full agentic loop using Groq as fallback."""
    messages = [{"role": "system", "content": system_prompt}]
    messages.extend(history)
    messages.append({"role": "user", "content": question})

    model = GROQ_MODEL

    for round_num in range(MAX_TOOL_ROUNDS):
        print(f"  Round {round_num+1}/{MAX_TOOL_ROUNDS} | Groq {model}")

        while True:
            try:
                resp = get_groq_client().chat.completions.create(
                    model=model,
                    messages=messages,
                    tools=GROQ_TOOLS,
                    tool_choice="auto",
                    parallel_tool_calls=False,
                    max_tokens=1500,
                    temperature=0.0,
                )
                break

            except Exception as e:
                err = str(e)
                if "429" in err or "rate_limit_exceeded" in err:
                    if rotate_groq():
                        time.sleep(2); continue
                    else:
                        raise  # all Groq keys exhausted
                elif any(c in err for c in ["model_decommissioned","model_not_found"]):
                    model = GROQ_MODEL_FAST
                    reset_groq()
                    break
                elif "tool_use_failed" in err:
                    # retry without tools
                    resp = get_groq_client().chat.completions.create(
                        model=model, messages=messages,
                        max_tokens=1500, temperature=0.0,
                    )
                    return (resp.choices[0].message.content or "").strip()
                else:
                    raise

        msg  = resp.choices[0].message
        stop = resp.choices[0].finish_reason
        print(f"  Stop: {stop}")

        if stop == "tool_calls" and msg.tool_calls:
            messages.append({
                "role": "assistant",
                "content": msg.content or "",
                "tool_calls": [
                    {"id": tc.id, "type": "function",
                     "function": {"name": tc.function.name,
                                  "arguments": tc.function.arguments}}
                    for tc in msg.tool_calls
                ]
            })
            for tc in msg.tool_calls:
                try:   args = json.loads(tc.function.arguments)
                except: args = {}
                result = dispatch(tc.function.name, args)
                messages.append({
                    "role": "tool",
                    "tool_call_id": tc.id,
                    "content": result
                })
        else:
            answer = (msg.content or "").strip()
            return answer or "I was unable to find a clear answer."

    return "I searched but could not compose a complete answer."


# ═══════════════════════════════════════════════════════════════════
# CHAT HISTORY  (static questions only)
# ═══════════════════════════════════════════════════════════════════
static_history = []
MAX_MEMORY     = 6

def update_history(q: str, a: str):
    static_history.append({"role": "user",      "content": q})
    static_history.append({"role": "assistant", "content": a})
    if len(static_history) > MAX_MEMORY * 2:
        static_history[:] = static_history[-(MAX_MEMORY*2):]


# ═══════════════════════════════════════════════════════════════════
# FILE LOADERS
# ═══════════════════════════════════════════════════════════════════
def load_pdf(path):
    try:   return [d.page_content for d in PyPDFLoader(path).load()]
    except Exception as e: print(f"  ⚠ PDF: {e}"); return []

def load_txt(path):
    with open(path,"r",encoding="utf-8") as f: return [f.read()]

def load_docx(path):
    return ["\n".join(p.text for p in Document(path).paragraphs)]

def load_pptx(path):
    prs, text = Presentation(path), []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"): text.append(shape.text)
    return ["\n".join(text)]

def load_file(path):
    ext = path.lower().rsplit(".",1)[-1]
    return {"pdf":load_pdf,"txt":load_txt,"docx":load_docx,
            "pptx":load_pptx}.get(ext, lambda p:[])(path)


# ═══════════════════════════════════════════════════════════════════
# FAISS INDEX
# ═══════════════════════════════════════════════════════════════════
print("Loading embedding model...")
embedding_model = SentenceTransformer("intfloat/multilingual-e5-small")

def build_index(paths):
    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
    all_docs = []
    for path in paths:
        for content in load_file(path):
            for chunk in splitter.split_text(translate_to_english(content)):
                all_docs.append({"text": chunk, "source": path})
    if not all_docs:
        return faiss.IndexFlatL2(embedding_model.get_sentence_embedding_dimension()), []
    texts = [d["text"] for d in all_docs]
    embs  = embedding_model.encode(texts, batch_size=32, show_progress_bar=True)
    arr   = np.array(embs).astype("float32")
    idx   = faiss.IndexFlatL2(arr.shape[1])
    idx.add(arr)
    return idx, all_docs

pdf_folder = "pdfs"
os.makedirs(pdf_folder, exist_ok=True)
paths = [
    os.path.join(pdf_folder, f)
    for f in os.listdir(pdf_folder)
    if f.lower().rsplit(".",1)[-1] in ("pdf","txt","docx","pptx")
    and os.path.isfile(os.path.join(pdf_folder, f))
]
for p in paths: print(f"  Document: {os.path.basename(p)}")

if os.path.exists(FAISS_FILE) and os.path.exists(DOCS_FILE):
    print("Loading FAISS index...")
    index = faiss.read_index(FAISS_FILE)
    with open(DOCS_FILE,"rb") as f: docs = pickle.load(f)
    print(f"  Loaded. Vectors: {index.ntotal}")
else:
    print("Building FAISS index...")
    index, docs = build_index(paths)
    faiss.write_index(index, FAISS_FILE)
    with open(DOCS_FILE,"wb") as f: pickle.dump(docs, f)
    print(f"  Built. Vectors: {index.ntotal}")


# ═══════════════════════════════════════════════════════════════════
# MAIN PIPELINE
# ═══════════════════════════════════════════════════════════════════
def rag_pipeline(question: str) -> str:
    """
    v6.0 pipeline:
    
    1. Classify question (LIVE or STATIC)
    2. Check answer cache for LIVE questions
    3. Try Gemini first (best multilingual, 1M tokens/day)
       - If Gemini key rotated and still fails → try Groq
    4. Try Groq as fallback (good English, 100k/day per key)
    5. Cache LIVE answers, store STATIC in history
    6. Translate if needed (Gemini: native, Groq: Google Translate)
    """
    answer_language = detect_preferred_language(question)
    english_q       = translate_to_english(question)
    force_govlk     = is_govlk_topic(english_q) or is_govlk_topic(question)
    force_web       = (not force_govlk) and (must_use_web(english_q) or must_use_web(question))
    live_q          = force_govlk or force_web

    print(f"\n{'═'*55}")
    print(f"  Lang        : {answer_language}")
    print(f"  Q           : {english_q}")
    print(f"  Type        : {'LIVE 🔴' if live_q else 'STATIC 📄'}")
    print(f"  Gov.lk first: {'YES 🏛️' if force_govlk else 'no'}")
    print(f"{'═'*55}")

    # Check cache
    if live_q:
        cached = get_cached(english_q)
        if cached:
            return smart_translate(cached, answer_language, used_gemini=True)

    # History for static questions only
    history = [] if live_q else static_history[:]

    # Prompt — tell Gemini to answer in user's language directly
    # (Gemini natively supports Sinhala & Tamil)
    llm_lang      = answer_language  # Gemini answers natively
    system_prompt = build_prompt(llm_lang, force_web=force_web, force_govlk=force_govlk)

    used_gemini = False
    raw_answer  = None

    # ── Try Gemini (PRIMARY) ──────────────────────────────────────
    if GEMINI_CLIENTS:
        reset_gemini()
        for attempt in range(len(GEMINI_CLIENTS)):
            try:
                raw_answer  = gemini_agentic(english_q, system_prompt, history, answer_language)
                used_gemini = True
                break
            except Exception as e:
                err = str(e)
                print(f"  ⚠ Gemini error: {err[:80]}")
                if "429" in err or "quota" in err.lower() or "rate" in err.lower():
                    if rotate_gemini():
                        time.sleep(2)
                        continue
                # Non-quota error or all keys exhausted
                print(f"  → Falling back to Groq")
                break

    # ── Try Groq (FALLBACK) ───────────────────────────────────────
    if raw_answer is None and GROQ_CLIENTS:
        reset_groq()
        # For Groq, ask in English and translate after
        groq_prompt = build_prompt("English", force_web=force_web, force_govlk=force_govlk)
        try:
            raw_answer  = groq_agentic(english_q, groq_prompt, history)
            used_gemini = False
            print(f"  ✓ Groq answer ({len(raw_answer)} chars)")
        except Exception as e:
            print(f"  ⚠ Groq also failed: {e}")

    # ── All providers failed ──────────────────────────────────────
    if raw_answer is None:
        raw_answer = (
            "All AI providers are currently at their limit. "
            "Gemini resets at midnight Pacific time. "
            "Groq resets at midnight UTC (5:30 AM Sri Lanka time). "
            "Please try again shortly."
        )
        used_gemini = False

    # ── Store / cache ─────────────────────────────────────────────
    if live_q:
        set_cached(english_q, raw_answer)
    else:
        update_history(english_q, raw_answer)

    # ── Translate ─────────────────────────────────────────────────
    return smart_translate(raw_answer, answer_language, used_gemini=used_gemini)


# ═══════════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    g_cap = len(GEMINI_CLIENTS) * 1000
    r_cap = len(GROQ_CLIENTS)   * 100
    print(f"\n  GNOS RAG Bot v6.0")
    print(f"  PRIMARY  : Gemini {GEMINI_PRIMARY_MODEL} × {len(GEMINI_CLIENTS)} keys ({g_cap}k/day)")
    print(f"  FALLBACK : Groq {GROQ_MODEL} × {len(GROQ_CLIENTS)} keys ({r_cap}k/day)")
    print(f"  TOTAL    : ~{g_cap+r_cap}k tokens/day FREE")
    print(f"  Today    : {get_today()}")
    print(f"  Docs     : {index.ntotal} vectors")
    print(f"  Commands : 'exit' | 'clear'\n")

    while True:
        q = input("\nAsk: ").strip()
        if not q:           continue
        if q.lower() == "exit": break
        if q.lower() == "clear":
            static_history.clear()
            _answer_cache.clear()
            print("  Cleared.")
            continue
        print(f"\n  Answer:\n{rag_pipeline(q)}\n")